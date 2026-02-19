"""
OneStream Manufacturing Accelerator - Pitch Deck Generator
Generates a professional PowerPoint presentation.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Brand colors
DARK_NAVY = RGBColor(0x0B, 0x1D, 0x3A)
ACCENT_BLUE = RGBColor(0x00, 0x6E, 0xC7)
LIGHT_BLUE = RGBColor(0x4D, 0xA8, 0xDA)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF7)
DARK_GRAY = RGBColor(0x34, 0x3A, 0x40)
MEDIUM_GRAY = RGBColor(0x6C, 0x75, 0x7D)
GREEN = RGBColor(0x28, 0xA7, 0x45)
ORANGE = RGBColor(0xFD, 0x7E, 0x14)
RED = RGBColor(0xDC, 0x35, 0x45)
GOLD = RGBColor(0xFF, 0xC1, 0x07)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def add_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape


def add_rounded_shape(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multi_text(slide, left, top, width, height, lines, font_size=16,
                   color=DARK_GRAY, line_spacing=1.5, font_name="Calibri"):
    """lines = list of (text, bold, color_override) tuples"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, bold, col = line_data, False, color
        else:
            text = line_data[0]
            bold = line_data[1] if len(line_data) > 1 else False
            col = line_data[2] if len(line_data) > 2 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = col
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1))
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=15,
                    color=DARK_GRAY, bullet_color=ACCENT_BLUE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        if isinstance(item, tuple):
            p.text = item[0]
            p.font.bold = item[1]
        else:
            p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = Pt(8)
        p.level = 0
    return txBox


def add_icon_card(slide, left, top, width, height, icon_text, title, subtitle,
                  bg_color=WHITE, icon_bg=ACCENT_BLUE, title_color=DARK_NAVY):
    card = add_rounded_shape(slide, left, top, width, height, bg_color)

    icon_size = Inches(0.6)
    icon_shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, left + Inches(0.3), top + Inches(0.3), icon_size, icon_size
    )
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = icon_bg
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = icon_text
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Segoe UI Emoji"

    add_text_box(slide, left + Inches(0.3), top + Inches(1.1), width - Inches(0.6), Inches(0.4),
                 title, font_size=14, color=title_color, bold=True)
    add_text_box(slide, left + Inches(0.3), top + Inches(1.5), width - Inches(0.6), Inches(1.0),
                 subtitle, font_size=11, color=MEDIUM_GRAY)
    return card


# ============================================================
# SLIDE 1: Title Slide
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_background(slide, DARK_NAVY)

# Accent bar at top
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

# Left accent line
add_shape(slide, Inches(1.2), Inches(2.2), Inches(0.08), Inches(1.0), ACCENT_BLUE)

add_text_box(slide, Inches(1.6), Inches(2.0), Inches(8), Inches(1.0),
             "OneStream XF", font_size=22, color=LIGHT_BLUE, bold=False)
add_text_box(slide, Inches(1.6), Inches(2.5), Inches(10), Inches(1.2),
             "Manufacturing Accelerator", font_size=48, color=WHITE, bold=True)
add_text_box(slide, Inches(1.6), Inches(3.8), Inches(9), Inches(0.8),
             "Pre-Built Enterprise CPM Solution for Global Manufacturing",
             font_size=22, color=LIGHT_BLUE)

# Bottom info bar
add_shape(slide, Inches(0), Inches(6.5), SLIDE_W, Inches(1.0), RGBColor(0x08, 0x15, 0x2B))
add_text_box(slide, Inches(1.6), Inches(6.65), Inches(4), Inches(0.6),
             "Confidential  |  2026", font_size=13, color=MEDIUM_GRAY)

# Key stats on right
for i, (num, label) in enumerate([("74", "Business Rules"), ("6", "Modules"), ("14", "Dimensions")]):
    x = Inches(9.0) + Inches(i * 1.5)
    add_text_box(slide, x, Inches(6.55), Inches(1.2), Inches(0.35),
                 num, font_size=24, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.9), Inches(1.2), Inches(0.35),
                 label, font_size=9, color=MEDIUM_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2: The Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "THE CHALLENGE", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "OneStream Implementations Are Expensive and Slow",
             font_size=36, color=DARK_NAVY, bold=True)

# Pain point cards
pain_points = [
    ("$500K-$1M+", "Average Project Cost",
     "Traditional implementations built from scratch for each client drive costs into 7 figures."),
    ("6-12 Months", "Typical Timeline",
     "Custom development, testing, and iteration cycles drag projects well beyond initial estimates."),
    ("40% Overrun", "Budget & Schedule Risk",
     "Complex consolidation logic, multi-source ETL, and manufacturing-specific calcs introduce scope creep."),
    ("Talent Scarcity", "Limited OneStream Developers",
     "The OneStream talent pool is small. Finding architects who understand both finance and VB.NET is hard."),
]

for i, (stat, title, desc) in enumerate(pain_points):
    x = Inches(0.8) + Inches(i * 3.05)
    y = Inches(2.2)
    card = add_rounded_shape(slide, x, y, Inches(2.85), Inches(3.5), LIGHT_GRAY)

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.25), Inches(0.6),
                 stat, font_size=28, color=RED, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.0), Inches(2.25), Inches(0.4),
                 title, font_size=15, color=DARK_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.5), Inches(2.25), Inches(1.8),
                 desc, font_size=12, color=MEDIUM_GRAY)

# Bottom line
add_shape(slide, Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.04), LIGHT_GRAY)
add_text_box(slide, Inches(0.8), Inches(6.4), Inches(11), Inches(0.5),
             "Manufacturing companies need a faster, lower-risk path to OneStream value.",
             font_size=15, color=DARK_NAVY, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 3: The Solution
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "THE SOLUTION", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "A Production-Ready Manufacturing Accelerator",
             font_size=36, color=DARK_NAVY, bold=True)

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(0.6),
             "Pre-built, source-controlled OneStream XF implementation covering the full CPM lifecycle for global multi-plant manufacturers.",
             font_size=16, color=MEDIUM_GRAY)

# Module cards (2 rows x 3)
modules = [
    ("FC", "Financial Consolidation",
     "Multi-currency, multi-GAAP consolidation with IC elimination, equity pickup, and minority interest",
     ACCENT_BLUE),
    ("PB", "Planning & Budgeting",
     "Driver-based planning with BOM rollups, headcount planning, CAPEX depreciation, rolling forecasts",
     RGBColor(0x6F, 0x42, 0xC1)),
    ("DM", "Data Management",
     "10 pre-built connectors for SAP, Oracle, NetSuite, Workday, MES systems with full ETL pipeline",
     GREEN),
    ("RD", "Reporting & Dashboards",
     "16 executive dashboards with variance waterfalls, KPI cockpits, and plant performance analytics",
     ORANGE),
    ("AR", "Account Reconciliation",
     "Automated matching engine with risk-based workflows and certification tracking",
     RGBColor(0xE8, 0x3E, 0x8C)),
    ("PP", "People Planning",
     "FTE-to-cost modeling: base salary, benefits, burden rates, merit increases, and org restructuring",
     LIGHT_BLUE),
]

for i, (icon, title, desc, color) in enumerate(modules):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(col * 4.0)
    y = Inches(2.6) + Inches(row * 2.3)

    card = add_rounded_shape(slide, x, y, Inches(3.8), Inches(2.1), LIGHT_GRAY)

    icon_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.25), y + Inches(0.25), Inches(0.55), Inches(0.55)
    )
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = color
    icon_shape.line.fill.background()
    tf = icon_shape.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Calibri"

    add_text_box(slide, x + Inches(1.0), y + Inches(0.25), Inches(2.6), Inches(0.4),
                 title, font_size=15, color=DARK_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(1.0), Inches(3.3), Inches(1.0),
                 desc, font_size=11, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 4: What's Included (By the Numbers)
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_NAVY)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "WHAT'S INCLUDED", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "194 Production-Ready Artifacts", font_size=36, color=WHITE, bold=True)

stats = [
    ("74", "VB.NET\nBusiness Rules", ACCENT_BLUE),
    ("14", "Dimension\nHierarchies", RGBColor(0x6F, 0x42, 0xC1)),
    ("27", "Data Mgmt\nPipelines", GREEN),
    ("16", "Executive\nDashboards", ORANGE),
    ("5", "Workflow\nDefinitions", RGBColor(0xE8, 0x3E, 0x8C)),
    ("10", "CubeView\nTemplates", LIGHT_BLUE),
    ("14", "Testing &\nValidation", GOLD),
    ("12", "Architecture\nDocs", RGBColor(0x20, 0xC9, 0x97)),
]

for i, (num, label, color) in enumerate(stats):
    col = i % 4
    row = i // 4
    x = Inches(0.8) + Inches(col * 3.1)
    y = Inches(2.0) + Inches(row * 2.6)

    card = add_rounded_shape(slide, x, y, Inches(2.8), Inches(2.3), RGBColor(0x11, 0x2B, 0x4A))

    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.8),
                 num, font_size=48, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(2.2), Inches(0.8),
                 label, font_size=14, color=RGBColor(0xAD, 0xB5, 0xBD))


# ============================================================
# SLIDE 5: Business Rules Deep Dive
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "BUSINESS RULES ENGINE", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "74 Production-Quality VB.NET Business Rules",
             font_size=36, color=DARK_NAVY, bold=True)

# Rule categories table
categories = [
    ("Finance Rules", "8", "Consolidation, FX translation, IC elimination, equity pickup, minority interest, goodwill, journal entries, flow analysis", ACCENT_BLUE),
    ("Calculate Rules", "20", "COGS allocation, overhead absorption, standard cost variance, OEE, revenue recognition, BOM rollup, driver-based planning, headcount, CAPEX, cash flow, KPIs", RGBColor(0x6F, 0x42, 0xC1)),
    ("Connector Rules", "10", "SAP HANA GL + production + materials, Oracle EBS GL + sub-ledger, NetSuite, Workday HCM, MES, Excel templates, flat files", GREEN),
    ("Dashboard Adapters", "15", "Executive summary, plant performance, variance waterfall, P&L bridge, BvA, rolling forecast, IC recon, CAPEX tracker, KPI cockpit", ORANGE),
    ("Member Filters", "5", "Entity security, scenario locking, time period control, product access, cost center access", LIGHT_BLUE),
    ("Event Handlers", "6", "Data quality validation, submission control, IC matching, budget threshold alerts, audit logging, notifications", RGBColor(0xE8, 0x3E, 0x8C)),
    ("Extenders", "6", "Batch consolidation, data archival, ETL orchestrator, recon engine, RF seeder, report distribution", GOLD),
    ("String Functions", "4", "Status icons, variance formatting, KPI thresholds, dynamic multi-language labels", RGBColor(0x20, 0xC9, 0x97)),
]

# Header
add_shape(slide, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.45), DARK_NAVY)
add_text_box(slide, Inches(1.0), Inches(1.83), Inches(2.5), Inches(0.4),
             "Category", font_size=12, color=WHITE, bold=True)
add_text_box(slide, Inches(3.5), Inches(1.83), Inches(0.8), Inches(0.4),
             "Count", font_size=12, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(1.83), Inches(7.5), Inches(0.4),
             "Key Capabilities", font_size=12, color=WHITE, bold=True)

for i, (cat, count, desc, color) in enumerate(categories):
    y = Inches(2.3) + Inches(i * 0.58)
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    add_shape(slide, Inches(0.8), y, Inches(11.7), Inches(0.55), bg)

    # Color indicator
    add_shape(slide, Inches(0.8), y, Inches(0.06), Inches(0.55), color)

    add_text_box(slide, Inches(1.0), y + Inches(0.08), Inches(2.5), Inches(0.4),
                 cat, font_size=12, color=DARK_NAVY, bold=True)
    add_text_box(slide, Inches(3.5), y + Inches(0.08), Inches(0.8), Inches(0.4),
                 count, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, Inches(4.5), y + Inches(0.08), Inches(7.8), Inches(0.4),
                 desc, font_size=10, color=MEDIUM_GRAY)

# Total bar
add_shape(slide, Inches(0.8), Inches(6.95), Inches(11.7), Inches(0.45), ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(6.98), Inches(2.5), Inches(0.4),
             "TOTAL", font_size=12, color=WHITE, bold=True)
add_text_box(slide, Inches(3.5), Inches(6.98), Inches(0.8), Inches(0.4),
             "74", font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.5), Inches(6.98), Inches(7.5), Inches(0.4),
             "53,000+ lines of production VB.NET code  |  All using native OneStream API patterns",
             font_size=10, color=WHITE, bold=True)


# ============================================================
# SLIDE 6: Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "ARCHITECTURE", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Enterprise-Grade Dimensional Model",
             font_size=36, color=DARK_NAVY, bold=True)

# Cube architecture
cubes = [
    ("Finance Cube", "12 Dimensions", "Consolidation, FX, IC elimination\n800+ accounts, 120 entities", ACCENT_BLUE),
    ("Planning Cube", "10 Dimensions", "Budget, forecast, driver-based\nplanning with version control", RGBColor(0x6F, 0x42, 0xC1)),
    ("HR Cube", "6 Dimensions", "People planning, comp modeling\nFTE-to-total-cost calculations", GREEN),
    ("Recon Cube", "4 Dimensions", "Account reconciliation with\nautomated matching engine", ORANGE),
]

for i, (name, dims, desc, color) in enumerate(cubes):
    x = Inches(0.8) + Inches(i * 3.1)
    card = add_rounded_shape(slide, x, y := Inches(1.8), Inches(2.85), Inches(1.8), LIGHT_GRAY)
    add_shape(slide, x, y, Inches(2.85), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.2), Inches(2.45), Inches(0.35),
                 name, font_size=16, color=DARK_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.55), Inches(2.45), Inches(0.3),
                 dims, font_size=12, color=color, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.9), Inches(2.45), Inches(0.8),
                 desc, font_size=10, color=MEDIUM_GRAY)

# Dimension list
add_text_box(slide, Inches(0.8), Inches(3.9), Inches(5), Inches(0.4),
             "14 Dimensions Fully Defined", font_size=16, color=DARK_NAVY, bold=True)

dims_left = [
    "Account  -  800+ members (P&L, BS, CF, Statistical)",
    "Entity  -  120 members (Corp > Region > Country > Plant)",
    "Scenario  -  27 members (Actual, Budget, RF01-12, WhatIf)",
    "Time  -  CY2024-2027 (Monthly, Quarterly, Annual)",
    "Flow  -  Opening, Movement, Closing, FX, Elimination",
    "Consolidation  -  Local, Translated, Proportional, Eliminated",
    "UD1 Product  -  Industrial / Consumer / Specialty SKUs",
]
dims_right = [
    "UD2 CostCenter  -  Production, Warehouse, Admin, R&D, Sales",
    "UD3 Intercompany  -  Mirrors Entity for IC matching",
    "UD4 Project  -  CAPEX Expansion, Maintenance, Automation, R&D",
    "UD5 CustomerSegment  -  OEM, Aftermarket, Govt, Commercial",
    "UD6 Channel  -  Direct, Distributor, OEM, E-Commerce",
    "UD7 Plant  -  Cross-ref for multi-entity plant analysis",
    "UD8 Version  -  Working, Submitted, Approved, Published",
]

for i, dim in enumerate(dims_left):
    add_text_box(slide, Inches(1.0), Inches(4.4) + Inches(i * 0.38), Inches(5.5), Inches(0.35),
                 dim, font_size=10, color=DARK_GRAY)

for i, dim in enumerate(dims_right):
    add_text_box(slide, Inches(7.0), Inches(4.4) + Inches(i * 0.38), Inches(5.5), Inches(0.35),
                 dim, font_size=10, color=DARK_GRAY)


# ============================================================
# SLIDE 7: Integration Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "DATA INTEGRATION", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Pre-Built Multi-Source ETL Pipeline",
             font_size=36, color=DARK_NAVY, bold=True)

# Source systems
sources = [
    ("SAP S/4HANA", "ODBC/RFC", "GL, Production Orders\nMaterial Master, BOM"),
    ("Oracle EBS", "DB Link", "General Ledger\nAP/AR Sub-Ledger"),
    ("NetSuite", "SuiteAnalytics", "GL Trial Balance\nRevenue Detail"),
    ("Workday", "REST API", "Headcount, Compensation\nOrg Hierarchy"),
    ("MES Systems", "SFTP/API", "Production Volumes\nOEE, Scrap, Quality"),
]

# Source column
add_text_box(slide, Inches(0.8), Inches(1.8), Inches(2.5), Inches(0.4),
             "Source Systems", font_size=14, color=DARK_NAVY, bold=True)

for i, (name, method, data) in enumerate(sources):
    y = Inches(2.3) + Inches(i * 0.95)
    card = add_rounded_shape(slide, Inches(0.8), y, Inches(2.8), Inches(0.85), LIGHT_GRAY)
    add_text_box(slide, Inches(1.0), y + Inches(0.05), Inches(2.4), Inches(0.3),
                 name, font_size=12, color=DARK_NAVY, bold=True)
    add_text_box(slide, Inches(1.0), y + Inches(0.35), Inches(1.0), Inches(0.2),
                 method, font_size=9, color=ACCENT_BLUE, bold=True)
    add_text_box(slide, Inches(2.0), y + Inches(0.35), Inches(1.6), Inches(0.45),
                 data, font_size=8, color=MEDIUM_GRAY)

# Pipeline stages
stages = [
    ("EXTRACT", "10 Connector\nBusiness Rules", ACCENT_BLUE),
    ("STAGE", "8 Staging\nDefinitions", RGBColor(0x6F, 0x42, 0xC1)),
    ("TRANSFORM", "8 Mapping &\nValidation Rules", GREEN),
    ("LOAD", "6 Target Cube\nLoad Sequences", ORANGE),
]

# Arrow pipeline
for i, (stage, desc, color) in enumerate(stages):
    x = Inches(4.2) + Inches(i * 2.3)
    y = Inches(3.2)

    card = add_rounded_shape(slide, x, y, Inches(2.0), Inches(2.0), color)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.3), Inches(1.7), Inches(0.4),
                 stage, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), y + Inches(0.9), Inches(1.7), Inches(0.8),
                 desc, font_size=11, color=RGBColor(0xFF, 0xFF, 0xEE), alignment=PP_ALIGN.CENTER)

    # Arrow between stages
    if i < len(stages) - 1:
        arrow_x = x + Inches(2.0)
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, y + Inches(0.75),
                                        Inches(0.3), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

# Bottom note
add_shape(slide, Inches(0.8), Inches(6.4), Inches(11.7), Inches(0.8), LIGHT_GRAY)
add_text_box(slide, Inches(1.0), Inches(6.5), Inches(11.3), Inches(0.6),
             "Full orchestration engine (EX_IntegrationOrchestrator.vb) coordinates multi-source extractions "
             "with retry logic, error handling, data quality scoring, and automated alerting.",
             font_size=12, color=DARK_GRAY)


# ============================================================
# SLIDE 8: Manufacturing-Specific Capabilities
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "MANUFACTURING FOCUS", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Built for the Complexity of Global Manufacturing",
             font_size=36, color=DARK_NAVY, bold=True)

capabilities = [
    ("Standard Cost Variance", "4-way decomposition: Price, Usage, Efficiency, Volume. "
     "Actual vs standard at every production line."),
    ("BOM Cost Rollup", "Recursive bill-of-materials explosion with memoization. "
     "Material + labor + overhead per unit."),
    ("OEE Calculation", "Availability x Performance x Quality. Machine-level and "
     "plant-level with drill-down."),
    ("WIP Valuation", "Materials issued + labor applied + overhead applied - completed goods. "
     "Real-time WIP tracking."),
    ("Capacity Utilization", "Actual machine hours vs available capacity. Bottleneck "
     "identification and shift analysis."),
    ("Transfer Pricing", "4 OECD methods: CUP, Cost Plus, Resale Price, TNMM. "
     "Arm's-length IC pricing engine."),
    ("P&L Lever Model", "8 levers: Volume, Price, Mix, FX, Input Cost, Labor, "
     "Overhead, SGA. Full waterfall decomposition."),
    ("Driver-Based Planning", "Units x Price x Mix x Seasonal Factors. Revenue "
     "through EBITDA with automatic KPI derivation."),
]

for i, (title, desc) in enumerate(capabilities):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(col * 6.2)
    y = Inches(1.8) + Inches(row * 1.3)

    add_shape(slide, x, y, Inches(0.06), Inches(1.1), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.25), y, Inches(5.6), Inches(0.35),
                 title, font_size=14, color=DARK_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.4), Inches(5.6), Inches(0.7),
                 desc, font_size=11, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 9: Time & Cost Savings
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_NAVY)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "VALUE PROPOSITION", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Accelerate Time-to-Value by 50%+",
             font_size=36, color=WHITE, bold=True)

# Comparison table
# Traditional column
add_rounded_shape(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(5.0),
                  RGBColor(0x11, 0x2B, 0x4A))
add_text_box(slide, Inches(1.2), Inches(2.2), Inches(4.7), Inches(0.5),
             "Traditional Implementation", font_size=18, color=RED, bold=True)

trad_items = [
    "Timeline: 6-12 months",
    "Cost: $500K - $1M+",
    "Team: 3-5 consultants",
    "Risk: High (custom build from scratch)",
    "Rules developed: 0 pre-built",
    "Dimensions: Designed from blank slate",
    "ETL: Connectors built per engagement",
    "Testing: Manual, ad-hoc validation",
]

for i, item in enumerate(trad_items):
    add_text_box(slide, Inches(1.5), Inches(2.9) + Inches(i * 0.48), Inches(4.5), Inches(0.4),
                 "x   " + item, font_size=13, color=RGBColor(0xAD, 0xB5, 0xBD))

# Accelerator column
add_rounded_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(5.0),
                  RGBColor(0x11, 0x2B, 0x4A))
add_shape(slide, Inches(7.0), Inches(2.0), Inches(5.5), Inches(0.06), GREEN)
add_text_box(slide, Inches(7.4), Inches(2.2), Inches(4.7), Inches(0.5),
             "With Manufacturing Accelerator", font_size=18, color=GREEN, bold=True)

accel_items = [
    "Timeline: 3-5 months",
    "Cost: $250K - $400K",
    "Team: 1-2 consultants",
    "Risk: Low (proven, tested framework)",
    "Rules included: 74 production-ready",
    "Dimensions: 14 pre-built hierarchies",
    "ETL: 10 connectors ready to configure",
    "Testing: 14 automated validation scripts",
]

for i, item in enumerate(accel_items):
    add_text_box(slide, Inches(7.3), Inches(2.9) + Inches(i * 0.48), Inches(4.5), Inches(0.4),
                 item, font_size=13, color=WHITE, bold=True if i < 2 else False)


# ============================================================
# SLIDE 10: Deployment & DevOps
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "DEVOPS & DEPLOYMENT", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Source-Controlled, CI/CD-Ready Deployment",
             font_size=36, color=DARK_NAVY, bold=True)

# Environment pipeline
envs = [
    ("DEV", "Development", "Debug logging, no rate limits\nPerformance profiling enabled\n20 concurrent users", ACCENT_BLUE),
    ("QA", "Quality Assurance", "Production-mirror data (masked)\nStrict validation, Info logging\n50 concurrent users", ORANGE),
    ("PROD", "Production", "HA with F5 load balancer\n99.9% SLA, MFA required\n250 concurrent users", GREEN),
]

for i, (code, name, desc, color) in enumerate(envs):
    x = Inches(0.8) + Inches(i * 4.2)
    card = add_rounded_shape(slide, x, y := Inches(1.8), Inches(3.8), Inches(2.5), LIGHT_GRAY)
    add_shape(slide, x, y, Inches(3.8), Inches(0.06), color)

    badge = add_rounded_shape(slide, x + Inches(0.2), y + Inches(0.2), Inches(0.8), Inches(0.4), color)
    tf = badge.text_frame
    tf.paragraphs[0].text = code
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Calibri"

    add_text_box(slide, x + Inches(1.2), y + Inches(0.2), Inches(2.4), Inches(0.35),
                 name, font_size=14, color=DARK_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(3.4), Inches(1.5),
                 desc, font_size=11, color=MEDIUM_GRAY)

    if i < 2:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                                        x + Inches(3.8), y + Inches(1.0), Inches(0.4), Inches(0.5))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = DARK_GRAY
        arrow.line.fill.background()

# Automation scripts
add_text_box(slide, Inches(0.8), Inches(4.6), Inches(5), Inches(0.4),
             "PowerShell Deployment Automation", font_size=16, color=DARK_NAVY, bold=True)

scripts = [
    ("Deploy_BusinessRules.ps1", "Uploads, compiles, and activates 74 business rules with rollback capability"),
    ("Deploy_Dimensions.ps1", "Loads dimension hierarchies with backup and delta comparison"),
    ("Deploy_DataManagement.ps1", "Deploys ETL sequences in dependency order with validation"),
    ("Validate_Deployment.ps1", "Post-deployment health checks: compilation, connectivity, data integrity"),
]

for i, (script, desc) in enumerate(scripts):
    y = Inches(5.1) + Inches(i * 0.52)
    add_text_box(slide, Inches(1.0), y, Inches(3.2), Inches(0.4),
                 script, font_size=11, color=ACCENT_BLUE, bold=True, font_name="Consolas")
    add_text_box(slide, Inches(4.5), y, Inches(8), Inches(0.4),
                 desc, font_size=11, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 11: ROI Analysis
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "RETURN ON INVESTMENT", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Quantifiable Business Impact",
             font_size=36, color=DARK_NAVY, bold=True)

# ROI metrics
roi_items = [
    ("50-60%", "Faster Implementation", "3-5 months vs 6-12 months. Faster time-to-value means earlier realization of CPM benefits and reduced project risk.", GREEN),
    ("$250-500K", "Cost Savings", "Reduced consulting hours, smaller team required, fewer custom development cycles. Accelerator pays for itself on the first engagement.", ACCENT_BLUE),
    ("80%", "Code Reusability", "74 business rules adapted, not rewritten. Dimension hierarchies customized, not designed from scratch. Proven patterns reduce defects.", RGBColor(0x6F, 0x42, 0xC1)),
    ("90%", "Risk Reduction", "Pre-tested consolidation logic, validated FX translation, proven IC elimination. Eliminates the #1 source of implementation failures.", ORANGE),
]

for i, (metric, title, desc, color) in enumerate(roi_items):
    y = Inches(1.9) + Inches(i * 1.3)
    card = add_rounded_shape(slide, Inches(0.8), y, Inches(11.7), Inches(1.15), LIGHT_GRAY)

    metric_box = add_rounded_shape(slide, Inches(1.0), y + Inches(0.15), Inches(1.8), Inches(0.85), color)
    tf = metric_box.text_frame
    tf.paragraphs[0].text = metric
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "Calibri"

    add_text_box(slide, Inches(3.1), y + Inches(0.1), Inches(9), Inches(0.35),
                 title, font_size=16, color=DARK_NAVY, bold=True)
    add_text_box(slide, Inches(3.1), y + Inches(0.5), Inches(9), Inches(0.6),
                 desc, font_size=12, color=MEDIUM_GRAY)


# ============================================================
# SLIDE 12: Engagement Model
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, WHITE)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.5),
             "ENGAGEMENT MODEL", font_size=13, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(0.8), Inches(0.8), Inches(11), Inches(0.8),
             "Flexible Delivery Options",
             font_size=36, color=DARK_NAVY, bold=True)

tiers = [
    ("Accelerator License", "$75K - $125K",
     [
         "Full source code repository access",
         "74 business rules + all artifacts",
         "Documentation & deployment scripts",
         "Self-service implementation",
         "Best for: Firms with in-house OneStream team",
     ],
     ACCENT_BLUE),
    ("Guided Implementation", "$200K - $350K",
     [
         "Everything in Accelerator License",
         "Solution architect engagement (8-12 weeks)",
         "Dimension & rule customization",
         "Integration configuration & testing",
         "Best for: Companies with some OneStream experience",
     ],
     RGBColor(0x6F, 0x42, 0xC1)),
    ("Full Delivery", "$350K - $600K",
     [
         "Everything in Guided Implementation",
         "Full project team (architect + 1-2 consultants)",
         "End-to-end delivery: design through go-live",
         "UAT support & hypercare period",
         "Best for: Companies new to OneStream",
     ],
     GREEN),
]

for i, (tier_name, price, features, color) in enumerate(tiers):
    x = Inches(0.8) + Inches(i * 4.1)
    card = add_rounded_shape(slide, x, Inches(1.8), Inches(3.8), Inches(5.2), LIGHT_GRAY)
    add_shape(slide, x, Inches(1.8), Inches(3.8), Inches(0.06), color)

    add_text_box(slide, x + Inches(0.3), Inches(2.1), Inches(3.2), Inches(0.4),
                 tier_name, font_size=16, color=DARK_NAVY, bold=True)
    add_text_box(slide, x + Inches(0.3), Inches(2.5), Inches(3.2), Inches(0.5),
                 price, font_size=24, color=color, bold=True)

    for j, feat in enumerate(features):
        is_last = j == len(features) - 1
        add_text_box(slide, x + Inches(0.3), Inches(3.2) + Inches(j * 0.55),
                     Inches(3.2), Inches(0.5),
                     feat, font_size=11,
                     color=color if is_last else DARK_GRAY,
                     bold=is_last)


# ============================================================
# SLIDE 13: Closing / CTA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide, DARK_NAVY)
add_shape(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ACCENT_BLUE)

# Left accent
add_shape(slide, Inches(1.2), Inches(2.2), Inches(0.08), Inches(1.0), ACCENT_BLUE)

add_text_box(slide, Inches(1.6), Inches(2.0), Inches(8), Inches(0.5),
             "Ready to Accelerate?", font_size=22, color=LIGHT_BLUE)
add_text_box(slide, Inches(1.6), Inches(2.5), Inches(10), Inches(1.0),
             "Let's Transform Your\nManufacturing Finance Operations",
             font_size=44, color=WHITE, bold=True)

add_text_box(slide, Inches(1.6), Inches(4.2), Inches(9), Inches(0.6),
             "From months of custom development to weeks of guided configuration.",
             font_size=20, color=LIGHT_BLUE)

# CTA boxes
ctas = [
    ("Schedule a Demo", "See the accelerator in action\nwith your data"),
    ("Technical Deep-Dive", "Architecture review with\nyour OneStream team"),
    ("Proof of Concept", "4-week pilot with your\nactual chart of accounts"),
]

for i, (title, desc) in enumerate(ctas):
    x = Inches(1.6) + Inches(i * 3.6)
    card = add_rounded_shape(slide, x, Inches(5.2), Inches(3.2), Inches(1.4),
                              RGBColor(0x11, 0x2B, 0x4A))
    add_shape(slide, x, Inches(5.2), Inches(3.2), Inches(0.05), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.25), Inches(5.4), Inches(2.7), Inches(0.4),
                 title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.25), Inches(5.8), Inches(2.7), Inches(0.6),
                 desc, font_size=11, color=RGBColor(0xAD, 0xB5, 0xBD))


# ============================================================
# SAVE
# ============================================================
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                           "OneStream_Manufacturing_Accelerator_PitchDeck.pptx")
prs.save(output_path)
print(f"Pitch deck saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
